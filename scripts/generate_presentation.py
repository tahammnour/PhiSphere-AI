"""
Generate PhiSphere AI Hackathon Presentation (.pptx)
Microsoft Innovation Challenge — March 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BG = RGBColor(0x04, 0x0B, 0x16)
NAVY = RGBColor(0x0D, 0x18, 0x29)
TEAL = RGBColor(0x00, 0xC9, 0xB1)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
INDIGO = RGBColor(0x81, 0x8C, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_GRAY = RGBColor(0x33, 0x40, 0x55)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
RED_SOFT = RGBColor(0xEF, 0x44, 0x44)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                     color=LIGHT_GRAY, bullet_color=TEAL, title=None, title_color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(font_size + 4)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.space_after = Pt(10)

    for i, item in enumerate(items):
        if title or i > 0:
            p = tf.add_paragraph()
        elif i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, title, body, accent_color=TEAL):
    add_shape_rect(slide, left, top, width, height, NAVY, DARK_GRAY)
    add_shape_rect(slide, left, top, Inches(0.06), height, accent_color)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.5),
                Inches(0.4), title, font_size=14, color=accent_color, bold=True)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.5), width - Inches(0.5),
                height - Inches(0.6), body, font_size=12, color=LIGHT_GRAY)


def add_table_slide(slide, left, top, rows_data, col_widths, header_bg=NAVY,
                    header_color=TEAL, body_color=LIGHT_GRAY):
    cols = len(col_widths)
    rows_count = len(rows_data)
    table_shape = slide.shapes.add_table(rows_count, cols, left, top,
                                         sum(col_widths), Inches(0.45) * rows_count)
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13 if ri == 0 else 12)
                paragraph.font.name = "Calibri"
                paragraph.font.bold = ri == 0
                paragraph.font.color.rgb = header_color if ri == 0 else body_color
                paragraph.alignment = PP_ALIGN.LEFT
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x0A, 0x12, 0x20) if ri == 0 else (
                NAVY if ri % 2 == 1 else RGBColor(0x0B, 0x15, 0x25)
            )
    return table_shape


# ─────────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), TEAL)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
            "PhiSphere AI", font_size=54, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
            "Lab Notebook AI Assistant", font_size=32, color=TEAL, bold=False,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
            "An AI lab partner for experimental reasoning and analysis",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.04), TEAL)

add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(0.5),
            "Microsoft Innovation Challenge \u2014 March 2026",
            font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
            "March 16\u201327, 2026",
            font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), TEAL)

# ─────────────────────────────────────────────
# SLIDE 2: The Problem
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "The Problem", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.05), TEAL)

problems = [
    "Researchers spend hours manually reasoning over experimental data, protocols, and results",
    "No unified tool for multimodal analysis \u2014 CSV data, lab images, PDF protocols are handled separately",
    "Existing AI tools lack safety boundaries for sensitive domains (biology, clinical, chemical)",
    "No transparency \u2014 AI suggestions come without explanation or audit trail",
    "Data quality issues (missing values, small samples) go undetected, leading to flawed conclusions",
    "No connection to public scientific datasets (OpenML, Kaggle) for benchmarking and comparison",
]

add_bullet_frame(slide, Inches(1.2), Inches(1.6), Inches(10.5), Inches(5.2),
                 [f"\u2022  {p}" for p in problems], font_size=18, color=LIGHT_GRAY)

# ─────────────────────────────────────────────
# SLIDE 3: Our Solution
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Our Solution: PhiSphere AI", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), TEAL)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
            "An agentic lab notebook assistant that helps researchers reason over experiments,\n"
            "analyze multimodal data, and safely suggest next steps \u2014 powered by Azure AI.",
            font_size=16, color=LIGHT_GRAY)

cards = [
    ("Structured AI Reasoning", "Every response uses:\nObservation \u2192 Analysis \u2192\nSuggested Next Steps \u2192\nWhy I Recommend This", TEAL),
    ("Multimodal Data", "CSV with auto-statistics & charts\nImages with Azure Vision\nPDFs with Document Intelligence\n+ RAG indexing", BLUE),
    ("Safety-First Design", "All AI responses fully buffered\nserver-side \u2192 Content Safety\nscreening \u2192 then streamed\nto client", EMERALD),
    ("Full Audit Trail", "Per-message safety metadata\nGroundedness scoring\nConfidence badges\nComplete audit log", VIOLET),
]

for i, (title, body, color) in enumerate(cards):
    col = i % 4
    add_card(slide, Inches(0.8 + col * 3.05), Inches(2.7), Inches(2.8), Inches(2.6),
             title, body, color)

features = [
    "12 protocol templates  \u2022  OpenML dataset import  \u2022  Hypothesis generation  \u2022  Draft paper export  \u2022  Metrics dashboard"
]
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5),
            features[0], font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 4: Architecture
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "System Architecture", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(2.0), Inches(0.05), TEAL)

box_h = Inches(0.7)

fe_left, fe_top = Inches(0.8), Inches(2.0)
add_shape_rect(slide, fe_left, fe_top, Inches(2.5), Inches(1.6), NAVY, TEAL)
add_textbox(slide, fe_left + Inches(0.15), fe_top + Inches(0.1), Inches(2.2), Inches(0.4),
            "FRONTEND", font_size=11, color=TEAL, bold=True)
add_textbox(slide, fe_left + Inches(0.15), fe_top + Inches(0.45), Inches(2.2), Inches(1.0),
            "React 19 + Vite 7\nTailwind CSS 4\nRecharts\nFramer Motion",
            font_size=12, color=LIGHT_GRAY)

be_left, be_top = Inches(4.2), Inches(2.0)
add_shape_rect(slide, be_left, be_top, Inches(2.5), Inches(1.6), NAVY, BLUE)
add_textbox(slide, be_left + Inches(0.15), be_top + Inches(0.1), Inches(2.2), Inches(0.4),
            "BACKEND", font_size=11, color=BLUE, bold=True)
add_textbox(slide, be_left + Inches(0.15), be_top + Inches(0.45), Inches(2.2), Inches(1.0),
            "Express 5 + Node 24\nPostgreSQL + Drizzle\nSSE Streaming\nRate Limiting",
            font_size=12, color=LIGHT_GRAY)

db_left, db_top = Inches(4.2), Inches(4.2)
add_shape_rect(slide, db_left, db_top, Inches(2.5), Inches(1.2), NAVY, AMBER)
add_textbox(slide, db_left + Inches(0.15), db_top + Inches(0.1), Inches(2.2), Inches(0.35),
            "DATABASE", font_size=11, color=AMBER, bold=True)
add_textbox(slide, db_left + Inches(0.15), db_top + Inches(0.4), Inches(2.2), Inches(0.7),
            "PostgreSQL\nDrizzle ORM\nlab_sessions, messages,\naudit_events",
            font_size=11, color=LIGHT_GRAY)

azure_services = [
    ("Azure OpenAI", "gpt-4o"),
    ("Content Safety", "Pre-stream filter"),
    ("AI Vision", "Image analysis"),
    ("Doc Intelligence", "PDF extraction"),
    ("AI Search", "RAG grounding"),
    ("AI Language", "Entity extraction"),
    ("App Insights", "Telemetry"),
    ("Azure ML", "Experiment tracking"),
]

azure_top = Inches(1.8)
az_left_start = Inches(7.6)
add_shape_rect(slide, az_left_start - Inches(0.15), azure_top - Inches(0.15),
               Inches(5.3), Inches(4.6), RGBColor(0x08, 0x10, 0x1E), BLUE)
add_textbox(slide, az_left_start, azure_top - Inches(0.05), Inches(5), Inches(0.35),
            "AZURE AI SERVICES (8)", font_size=12, color=BLUE, bold=True)

for i, (name, desc) in enumerate(azure_services):
    row = i // 2
    col = i % 2
    x = az_left_start + Inches(0.15) + col * Inches(2.5)
    y = azure_top + Inches(0.4) + row * Inches(0.95)
    add_shape_rect(slide, x, y, Inches(2.3), Inches(0.8), NAVY, DARK_GRAY)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(2.0), Inches(0.3),
                name, font_size=11, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.4), Inches(2.0), Inches(0.3),
                desc, font_size=10, color=MID_GRAY)

ext_left, ext_top = Inches(0.8), Inches(4.2)
add_shape_rect(slide, ext_left, ext_top, Inches(2.5), Inches(1.2), NAVY, INDIGO)
add_textbox(slide, ext_left + Inches(0.15), ext_top + Inches(0.1), Inches(2.2), Inches(0.35),
            "EXTERNAL DATA", font_size=11, color=INDIGO, bold=True)
add_textbox(slide, ext_left + Inches(0.15), ext_top + Inches(0.4), Inches(2.2), Inches(0.7),
            "OpenML API (live import)\nKaggle Sensor Datasets\nSample datasets (built-in)",
            font_size=11, color=LIGHT_GRAY)

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
            "Monorepo: pnpm workspaces  \u2022  TypeScript 5.9  \u2022  OpenAPI 3.1 + Orval codegen  \u2022  Zod validation",
            font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 5: AI Pipeline (Data Flow)
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "AI Pipeline \u2014 Safety-First Data Flow", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(3.0), Inches(0.05), TEAL)

pipeline_steps = [
    ("1", "User Message", "Researcher sends question\nwith session context", TEAL),
    ("2", "RAG Search", "Azure AI Search retrieves\ntop-3 document chunks", BLUE),
    ("3", "Entity Extraction", "Azure Language extracts\nscientific entities", INDIGO),
    ("4", "LLM Reasoning", "Azure OpenAI generates\nstructured response", WHITE),
    ("5", "Safety Buffer", "Full response buffered\n\u2192 Content Safety check", AMBER),
    ("6", "Stream to Client", "SSE stream with RAG chips,\nentity pills, RAI panel", EMERALD),
]

y_base = Inches(1.8)
for i, (num, title, desc, color) in enumerate(pipeline_steps):
    x = Inches(0.5) + i * Inches(2.1)
    add_shape_rect(slide, x, y_base, Inches(1.9), Inches(2.2), NAVY, color)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y_base + Inches(0.15),
                                    Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BG
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.1), y_base + Inches(0.75), Inches(1.7), Inches(0.4),
                title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y_base + Inches(1.15), Inches(1.7), Inches(0.9),
                desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
            "KEY: All AI responses are fully buffered server-side before Content Safety runs. "
            "Only after safety passes does the response stream to the client.",
            font_size=14, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

key_features = [
    ("Groundedness Scoring", "Azure AI evaluates whether responses\nare grounded in provided data", TEAL),
    ("Confidence Badges", "Per-message High / Medium / Low\nconfidence visible to researchers", BLUE),
    ("Audit Trail", "Every safety event recorded with\ntimestamps and category details", EMERALD),
]

for i, (title, desc, color) in enumerate(key_features):
    x = Inches(0.8) + i * Inches(4.0)
    add_card(slide, x, Inches(5.2), Inches(3.7), Inches(1.5), title, desc, color)

# ─────────────────────────────────────────────
# SLIDE 6: Azure Services Deep Dive
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Azure Services \u2014 Breadth & Depth", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), TEAL)

azure_table_data = [
    ["#", "Azure Service", "Purpose", "How It's Used"],
    ["1", "Azure OpenAI", "GPT-4o chat completions", "Structured scientific reasoning with domain safety rules"],
    ["2", "AI Content Safety", "Pre-stream safety screening", "Full response buffer \u2192 analyze \u2192 pass/flag/block"],
    ["3", "AI Vision", "Image analysis", "Captions, OCR, object detection for lab photos"],
    ["4", "Document Intelligence", "PDF extraction", "Prebuilt-layout model \u2192 text chunks for RAG"],
    ["5", "AI Search", "RAG grounding", "Auto-create index, top-3 chunk retrieval with citations"],
    ["6", "AI Language", "Entity recognition", "Scientific terms, chemicals, genes at \u226570% confidence"],
    ["7", "Application Insights", "Telemetry & monitoring", "Custom events: SessionCreated, SafetyFlagTriggered, etc."],
    ["8", "Azure ML", "Experiment tracking", "MLflow logging, dataset registration, model metrics"],
]

add_table_slide(slide, Inches(0.5), Inches(1.6), azure_table_data,
                [Inches(0.4), Inches(2.2), Inches(2.5), Inches(7.5)])

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
            "Live health checks for all services via /api/azure/status  \u2022  "
            "Status panel shows active count badge (e.g. 7/8)",
            font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 7: Innovation Highlights
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Innovation Highlights", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(2.0), Inches(0.05), TEAL)

innovations = [
    ("OpenML Live Import", "Search and import any public dataset by ID\ndirectly from OpenML into a lab session.\nARFF \u2192 CSV auto-conversion included.", TEAL),
    ("Structured Reasoning", "Every AI response follows:\nObservation \u2192 Analysis \u2192 Next Steps\n\u2192 Why I Recommend This", BLUE),
    ("12 Protocol Templates", "PCR, Western Blot, ELISA, Gel Electro-\nphoresis, Cell Culture, Titration,\nSpectrophotometry, DNA Extraction...", INDIGO),
    ("Hypothesis Generation", "AI generates testable hypotheses with\nconfidence levels (high/medium/low)\nand reasoning chains", VIOLET),
    ("Draft Paper Export", "One-click research paper draft:\nAbstract, Materials & Methods,\nfull paper outline", EMERALD),
    ("Platform Metrics", "Real-time KPIs: safety pass rate,\ngroundedness scores, RAG usage,\nAzure service health", AMBER),
]

for i, (title, desc, color) in enumerate(innovations):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + col * Inches(4.15)
    y = Inches(1.6) + row * Inches(2.6)
    add_card(slide, x, y, Inches(3.9), Inches(2.2), title, desc, color)

# ─────────────────────────────────────────────
# SLIDE 8: Responsible AI
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Responsible AI \u2014 Safety at Every Layer", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(3.0), Inches(0.05), EMERALD)

rai_items = [
    ("Safety-First Streaming", "AI responses fully buffered server-side.\n"
     "Azure Content Safety screens complete text\nbefore any content reaches the client.", EMERALD),
    ("Domain Safety Rules", "Hard-coded system prompt boundaries:\n"
     "\u2022 No pathogen synthesis routes\n"
     "\u2022 No clinical dosing advice\n"
     "\u2022 No explosive synthesis\n"
     "\u2022 Mandatory BSL/IRB flags", RED_SOFT),
    ("Data Quality Warnings", "Automatic detection of:\n"
     "\u2022 Small sample sizes (<30 rows)\n"
     "\u2022 Missing values\n"
     "\u2022 Zero-variance columns\n"
     "\u2022 Low numeric column ratio", AMBER),
    ("Groundedness Scoring", "Azure AI evaluates whether responses\n"
     "are grounded in provided data vs.\n"
     "fabricated. Score stored per-message.", TEAL),
    ("Complete Audit Trail", "Every safety event persisted with:\n"
     "\u2022 Timestamps & severity categories\n"
     "\u2022 Session & conversation context\n"
     "\u2022 Model deployment metadata", BLUE),
    ("RAI Toolbox Notebooks", "Offline analysis using Microsoft's\n"
     "Responsible AI Toolbox for data balance,\n"
     "error analysis, and correlation checks.", VIOLET),
]

for i, (title, desc, color) in enumerate(rai_items):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + col * Inches(4.15)
    y = Inches(1.6) + row * Inches(2.7)
    add_card(slide, x, y, Inches(3.9), Inches(2.3), title, desc, color)

# ─────────────────────────────────────────────
# SLIDE 9: Judging Criteria Alignment
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Judging Criteria Alignment", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), TEAL)

criteria = [
    ("Performance (25%)",
     "\u2022 SSE streaming for sub-second first-token latency\n"
     "\u2022 Rate limiting (auth: 20/15min, AI: 20/min, general: 120/min)\n"
     "\u2022 Live health checks across all 8 Azure services\n"
     "\u2022 Metrics dashboard with real-time KPIs\n"
     "\u2022 esbuild for fast server bundling", TEAL),
    ("Innovation (25%)",
     "\u2022 Agentic reasoning pipeline (ingest \u2192 RAG \u2192 LLM \u2192 safety \u2192 audit)\n"
     "\u2022 OpenML live dataset import with ARFF conversion\n"
     "\u2022 Multimodal analysis: CSV + images + PDFs\n"
     "\u2022 12 protocol templates with starter prompts\n"
     "\u2022 Hypothesis generation + draft paper export", BLUE),
    ("Azure Breadth (25%)",
     "\u2022 8 Azure services integrated:\n"
     "  OpenAI, Content Safety, Vision, Document Intelligence,\n"
     "  AI Search, AI Language, App Insights, Azure ML\n"
     "\u2022 Each service used meaningfully, not superficially\n"
     "\u2022 Health monitoring for every service", INDIGO),
    ("Responsible AI (25%)",
     "\u2022 Server-side safety buffer before streaming\n"
     "\u2022 Domain-specific safety rules (bio, clinical, chemical)\n"
     "\u2022 Groundedness scoring + confidence badges\n"
     "\u2022 Data quality warnings (RAI)\n"
     "\u2022 Audit log + RAI Toolbox notebooks", EMERALD),
]

for i, (title, body, color) in enumerate(criteria):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.2)
    y = Inches(1.5) + row * Inches(2.8)
    add_card(slide, x, y, Inches(5.9), Inches(2.5), title, body, color)

# ─────────────────────────────────────────────
# SLIDE 10: Demo Walkthrough
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Demo Walkthrough", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(1.8), Inches(0.05), TEAL)

demo_steps = [
    "Landing page \u2192 Sign up \u2192 Control Panel with Azure service status",
    "Create session: Biology domain + PCR Amplification template",
    "Load sample CSV \u2192 auto-statistics + interactive charts",
    "Ask AI: \"What trends do you see in temperature vs. growth?\"",
    "Show structured response: Observation / Analysis / Next Steps / Why",
    "Import OpenML dataset #61 (Iris) \u2192 live external data",
    "Upload gel electrophoresis image \u2192 Azure Vision analysis",
    "Upload PDF protocol \u2192 RAG indexing + citation chips in chat",
    "Safety demo: borderline question \u2192 Content Safety blocks + audit log",
    "Responsible AI panel: safety check, groundedness, reasoning trace",
    "Platform metrics dashboard: safety pass rate, RAG usage, groundedness",
    "Export session as markdown lab notebook",
]

for i, step in enumerate(demo_steps):
    y = Inches(1.5) + i * Inches(0.46)
    num_color = TEAL if i < 4 else BLUE if i < 8 else EMERALD

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.04),
                                    Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = num_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BG
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(1.5), y, Inches(10.5), Inches(0.4),
                step, font_size=15, color=LIGHT_GRAY)

# ─────────────────────────────────────────────
# SLIDE 11: Tech Stack
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Technology Stack", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(2.0), Inches(0.05), TEAL)

tech_table = [
    ["Layer", "Technology", "Details"],
    ["Frontend", "React 19 + Vite 7 + Tailwind CSS 4", "Framer Motion, Recharts, Lucide, Wouter"],
    ["Backend", "Express 5 + Node.js 24", "SSE streaming, rate limiting, multer uploads"],
    ["Database", "PostgreSQL + Drizzle ORM", "lab_sessions, conversations, messages, audit_events"],
    ["API", "OpenAPI 3.1 + Orval + Zod", "Type-safe codegen, React Query hooks"],
    ["AI Engine", "Azure OpenAI (gpt-4o)", "Structured reasoning with domain safety rules"],
    ["Safety", "Azure AI Content Safety", "Server-side buffer, 4-category analysis"],
    ["Vision", "Azure AI Vision", "Captions, OCR, object detection"],
    ["RAG", "Azure AI Search + Doc Intelligence", "PDF \u2192 chunks \u2192 index \u2192 retrieval"],
    ["NLP", "Azure AI Language", "Named entity recognition (chemicals, genes)"],
    ["Telemetry", "Azure Application Insights", "Custom events, exception tracking"],
    ["MLOps", "Azure ML + MLflow", "Experiment tracking, dataset registration"],
    ["Build", "pnpm workspaces + esbuild + TypeScript 5.9", "Monorepo with shared packages"],
]

add_table_slide(slide, Inches(0.5), Inches(1.5), tech_table,
                [Inches(1.3), Inches(4.5), Inches(6.5)])

# ─────────────────────────────────────────────
# SLIDE 12: Key Learnings
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Key Learnings", font_size=36, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(1.8), Inches(0.05), TEAL)

learnings = [
    ("Safety Must Be Architecture, Not Afterthought",
     "Buffering all AI responses server-side before safety screening is essential.\n"
     "Streaming directly to the client makes it impossible to retract unsafe content.",
     EMERALD),
    ("RAG Grounding Dramatically Improves Quality",
     "Injecting document chunks into the LLM context with citations produces\n"
     "far more accurate, trustworthy responses than prompt-only approaches.",
     BLUE),
    ("Data Quality Warnings Prevent Bad Science",
     "Automatically flagging small samples, missing values, and zero-variance\n"
     "columns prevents researchers from drawing flawed conclusions.",
     AMBER),
    ("Azure Service Integration Requires Graceful Degradation",
     "Each service can fail independently. The system must degrade gracefully\n"
     "with clear status indicators, not crash when one API key is missing.",
     VIOLET),
]

for i, (title, body, color) in enumerate(learnings):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.2)
    y = Inches(1.5) + row * Inches(2.8)
    add_card(slide, x, y, Inches(5.9), Inches(2.4), title, body, color)

# ─────────────────────────────────────────────
# SLIDE 13: Thank You / Closing
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), TEAL)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
            "Thank You", font_size=54, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "PhiSphere AI \u2014 Lab Notebook AI Assistant",
            font_size=24, color=TEAL, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(5.5), Inches(3.9), Inches(2.3), Inches(0.04), TEAL)

links = [
    "GitHub:  github.com/your-repo/PhiSphere-AI",
    "Live Demo:  phisphere-ai.replit.app",
    "Microsoft Innovation Challenge \u2014 March 2026",
]
for i, link in enumerate(links):
    add_textbox(slide, Inches(1), Inches(4.4 + i * 0.5), Inches(11), Inches(0.5),
                link, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
            "Powered by Azure OpenAI \u2022 AI Vision \u2022 Content Safety \u2022 AI Search \u2022 "
            "Doc Intelligence \u2022 AI Language \u2022 App Insights \u2022 Azure ML",
            font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), TEAL)


# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
output_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
output_path = os.path.join(output_dir, "PhiSphere_AI_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
